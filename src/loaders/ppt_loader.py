from pptx import Presentation

def load_pptx(path):
    """
    Extracts text from a PowerPoint (.pptx) file.
    """
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
